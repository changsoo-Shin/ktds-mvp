"""
문서 처리 모듈
PDF, Word, 텍스트, Excel, PowerPoint 파일에서 내용을 추출하는 기능 제공
"""

import io
import os
from typing import Union, List, Dict
import PyPDF2
import docx
import pandas as pd
from pptx import Presentation
import olefile
import nbformat
from streamlit.runtime.uploaded_file_manager import UploadedFile

class DocumentProcessor:
    """문서 처리 클래스"""
    
    def __init__(self):
        self.supported_formats = ['.pdf', '.docx', '.doc', '.txt', '.xlsx', '.ppt', '.pptx', '.hwp', '.ipynb']
    
    def extract_content(self, file: UploadedFile) -> str:
        """
        업로드된 파일에서 텍스트 내용 추출
        
        Args:
            file: Streamlit UploadedFile 객체
            
        Returns:
            str: 추출된 텍스트 내용
        """
        file_extension = os.path.splitext(file.name)[1].lower()
        
        try:
            if file_extension == '.pdf':
                return self._extract_from_pdf(file)
            elif file_extension in ['.docx', '.doc']:
                return self._extract_from_word(file)
            elif file_extension == '.txt':
                return self._extract_from_text(file)
            elif file_extension == '.xlsx':
                return self._extract_from_excel(file)
            elif file_extension in ['.ppt', '.pptx']:
                return self._extract_from_powerpoint(file)
            elif file_extension == '.hwp':
                return self._extract_from_hwp(file)
            elif file_extension == '.ipynb':
                return self._extract_from_notebook(file)
            else:
                raise ValueError(f"지원하지 않는 파일 형식: {file_extension}")
                
        except Exception as e:
            raise Exception(f"문서 처리 중 오류 발생: {str(e)}")
    
    def _extract_from_pdf(self, file: UploadedFile) -> str:
        """PDF 파일에서 텍스트 추출"""
        try:
            # 파일을 바이트로 읽기
            file_bytes = file.read()
            
            # PDF 리더 생성
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
            
            text_content = ""
            
            # 모든 페이지에서 텍스트 추출
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text_content += page.extract_text() + "\n"
            
            if not text_content.strip():
                raise Exception("PDF에서 텍스트를 추출할 수 없습니다.")
            
            return text_content.strip()
            
        except Exception as e:
            raise Exception(f"PDF 처리 실패: {str(e)}")
    
    def _extract_from_word(self, file: UploadedFile) -> str:
        """Word 파일에서 텍스트 추출"""
        try:
            # 파일을 바이트로 읽기
            file_bytes = file.read()
            
            # Word 문서 열기
            doc = docx.Document(io.BytesIO(file_bytes))
            
            text_content = ""
            
            # 모든 단락에서 텍스트 추출
            for paragraph in doc.paragraphs:
                text_content += paragraph.text + "\n"
            
            # 표에서 텍스트 추출
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text_content += cell.text + " "
                    text_content += "\n"
            
            if not text_content.strip():
                raise Exception("Word 문서에서 텍스트를 추출할 수 없습니다.")
            
            return text_content.strip()
            
        except Exception as e:
            raise Exception(f"Word 문서 처리 실패: {str(e)}")
    
    def _extract_from_text(self, file: UploadedFile) -> str:
        """텍스트 파일에서 내용 읽기"""
        try:
            # 파일을 문자열로 디코딩
            file_bytes = file.read()
            
            # 다양한 인코딩 시도
            encodings = ['utf-8', 'cp949', 'euc-kr', 'latin-1']
            
            for encoding in encodings:
                try:
                    text_content = file_bytes.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            else:
                raise Exception("파일 인코딩을 인식할 수 없습니다.")
            
            if not text_content.strip():
                raise Exception("텍스트 파일이 비어있습니다.")
            
            return text_content.strip()
            
        except Exception as e:
            raise Exception(f"텍스트 파일 처리 실패: {str(e)}")
    
    def _extract_from_excel(self, file: UploadedFile) -> str:
        """Excel 파일에서 텍스트 추출"""
        try:
            # 파일을 바이트로 읽기
            file_bytes = file.read()
            
            # Excel 파일 읽기
            df_dict = pd.read_excel(io.BytesIO(file_bytes), sheet_name=None)
            
            text_content = ""
            
            # 모든 시트에서 데이터 추출
            for sheet_name, df in df_dict.items():
                text_content += f"\n=== {sheet_name} ===\n"
                
                # 컬럼명 추가
                if not df.empty:
                    text_content += "컬럼: " + ", ".join(df.columns.astype(str)) + "\n\n"
                    
                    # 각 행의 데이터를 텍스트로 변환
                    for index, row in df.iterrows():
                        row_text = []
                        for col, value in row.items():
                            if pd.notna(value):  # NaN이 아닌 값만 포함
                                row_text.append(f"{col}: {str(value)}")
                        
                        if row_text:  # 빈 행이 아닌 경우만 추가
                            text_content += " | ".join(row_text) + "\n"
                    
                    text_content += "\n"
            
            if not text_content.strip():
                raise Exception("Excel 파일에서 텍스트를 추출할 수 없습니다.")
            
            return text_content.strip()
            
        except Exception as e:
            raise Exception(f"Excel 파일 처리 실패: {str(e)}")
    
    def _extract_from_powerpoint(self, file: UploadedFile) -> str:
        """PowerPoint 파일에서 텍스트 추출"""
        try:
            # 파일을 바이트로 읽기
            file_bytes = file.read()
            
            # PowerPoint 파일 열기
            prs = Presentation(io.BytesIO(file_bytes))
            
            text_content = ""
            
            # 모든 슬라이드에서 텍스트 추출
            for slide_num, slide in enumerate(prs.slides, 1):
                text_content += f"\n=== 슬라이드 {slide_num} ===\n"
                
                # 슬라이드의 모든 shape에서 텍스트 추출
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content += shape.text.strip() + "\n"
                    
                    # 표가 있는 경우 표 내용도 추출
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                text_content += " | ".join(row_text) + "\n"
                
                text_content += "\n"
            
            if not text_content.strip():
                raise Exception("PowerPoint 파일에서 텍스트를 추출할 수 없습니다.")
            
            return text_content.strip()
            
        except Exception as e:
            raise Exception(f"PowerPoint 파일 처리 실패: {str(e)}")
    
    def _extract_from_hwp(self, file: UploadedFile) -> str:
        """한글 파일에서 텍스트 추출"""
        try:
            # 파일을 바이트로 읽기
            file_bytes = file.read()
            
            # 임시 파일로 저장 (olefile은 파일 경로가 필요함)
            import tempfile
            with tempfile.NamedTemporaryFile(delete=False, suffix='.hwp') as temp_file:
                temp_file.write(file_bytes)
                temp_file_path = temp_file.name
            
            try:
                # HWP 파일이 OLE 파일인지 확인
                if not olefile.isOleFile(temp_file_path):
                    raise Exception("유효하지 않은 HWP 파일입니다.")
                
                # OLE 파일 열기
                ole = olefile.OleFileIO(temp_file_path)
                
                text_content = ""
                
                # HWP 파일의 텍스트 스트림 찾기
                # HWP 5.0 이상의 경우 'BodyText' 디렉토리에서 섹션별로 텍스트 추출
                if ole._olestream_size:
                    try:
                        # BodyText 디렉토리 내의 섹션들을 찾아서 텍스트 추출
                        for stream_name in ole.listdir():
                            if isinstance(stream_name, list) and len(stream_name) > 0:
                                if stream_name[0] == 'BodyText':
                                    # BodyText 하위의 Section 스트림들 처리
                                    for section_stream in ole.listdir():
                                        if (isinstance(section_stream, list) and 
                                            len(section_stream) >= 2 and 
                                            section_stream[0] == 'BodyText' and 
                                            section_stream[1].startswith('Section')):
                                            
                                            try:
                                                # 섹션 스트림 읽기
                                                stream_data = ole._olestream(section_stream)
                                                
                                                # 간단한 텍스트 추출 (완전하지 않지만 기본적인 텍스트는 추출 가능)
                                                # HWP의 복잡한 구조로 인해 완벽한 텍스트 추출은 어려움
                                                decoded_text = self._extract_text_from_hwp_stream(stream_data)
                                                if decoded_text:
                                                    text_content += decoded_text + "\n"
                                                    
                                            except Exception as stream_error:
                                                # 개별 스트림 처리 실패는 무시하고 계속 진행
                                                continue
                                
                        # 만약 BodyText에서 추출이 안 되면 다른 방법 시도
                        if not text_content.strip():
                            # PrvText 스트림에서 텍스트 추출 시도
                            try:
                                if ole.exists('PrvText'):
                                    stream_data = ole._olestream(['PrvText'])
                                    decoded_text = self._extract_text_from_hwp_stream(stream_data)
                                    if decoded_text:
                                        text_content = decoded_text
                            except:
                                pass
                                
                    except Exception as extract_error:
                        # 텍스트 추출 실패 시 기본 메시지
                        text_content = f"한글 파일에서 텍스트를 추출하는 중 오류가 발생했습니다: {str(extract_error)}\n"
                        text_content += "한글 파일의 복잡한 구조로 인해 완전한 텍스트 추출이 어려울 수 있습니다."
                
                ole.close()
                
                if not text_content.strip():
                    # 대안적 방법: 바이너리 데이터에서 한글 텍스트 패턴 찾기
                    text_content = self._extract_korean_text_from_binary(file_bytes)
                    if not text_content.strip():
                        raise Exception("한글 파일에서 텍스트를 추출할 수 없습니다. 파일이 손상되었거나 지원되지 않는 형식일 수 있습니다.")
                
                return text_content.strip()
                
            finally:
                # 임시 파일 삭제
                try:
                    os.unlink(temp_file_path)
                except:
                    pass
                    
        except Exception as e:
            raise Exception(f"한글 파일 처리 실패: {str(e)}")
    
    def _extract_text_from_hwp_stream(self, stream_data: bytes) -> str:
        """HWP 스트림에서 텍스트 추출"""
        try:
            # HWP 파일의 텍스트는 UTF-16LE로 인코딩되어 있을 수 있음
            # 여러 인코딩 방식 시도
            encodings = ['utf-16le', 'utf-8', 'cp949', 'euc-kr']
            
            for encoding in encodings:
                try:
                    # 바이너리 데이터에서 텍스트 부분만 추출
                    text = stream_data.decode(encoding, errors='ignore')
                    # 제어 문자 및 불필요한 문자 제거
                    cleaned_text = ''.join(char for char in text if char.isprintable() or char.isspace())
                    # 한글이 포함된 경우만 반환
                    if any('\uac00' <= char <= '\ud7af' for char in cleaned_text):
                        return cleaned_text
                except:
                    continue
            
            return ""
            
        except Exception:
            return ""
    
    def _extract_korean_text_from_binary(self, file_bytes: bytes) -> str:
        """바이너리 데이터에서 한글 텍스트 패턴 추출"""
        try:
            # UTF-16LE로 디코딩 시도
            try:
                text = file_bytes.decode('utf-16le', errors='ignore')
                # 한글 문자만 추출
                korean_chars = []
                for char in text:
                    if '\uac00' <= char <= '\ud7af' or char.isspace() or char in '.,!?;:()[]{}':
                        korean_chars.append(char)
                
                korean_text = ''.join(korean_chars)
                # 연속된 공백 제거
                import re
                korean_text = re.sub(r'\s+', ' ', korean_text).strip()
                
                if len(korean_text) > 10:  # 최소 길이 확인
                    return korean_text
                    
            except:
                pass
            
            # CP949로 디코딩 시도
            try:
                text = file_bytes.decode('cp949', errors='ignore')
                # 한글이 포함된 텍스트 추출
                if any('\uac00' <= char <= '\ud7af' for char in text):
                    return text
            except:
                pass
            
            return ""
            
        except Exception:
            return ""
    
    def chunk_text(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """
        텍스트를 청크로 분할
        
        Args:
            text: 분할할 텍스트
            chunk_size: 청크 크기 (문자 수)
            overlap: 청크 간 겹치는 부분 (문자 수)
            
        Returns:
            List[str]: 분할된 텍스트 청크 리스트
        """
        if len(text) <= chunk_size:
            return [text]
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            
            # 문장 경계에서 분할 시도
            if end < len(text):
                # 마지막 문장 끝 찾기
                last_period = text.rfind('.', start, end)
                last_newline = text.rfind('\n', start, end)
                
                if last_period > start:
                    end = last_period + 1
                elif last_newline > start:
                    end = last_newline + 1
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            start = end - overlap
            if start >= len(text):
                break
        
        return chunks
    
    def extract_metadata(self, file: UploadedFile) -> Dict[str, str]:
        """
        파일 메타데이터 추출
        
        Args:
            file: Streamlit UploadedFile 객체
            
        Returns:
            Dict[str, str]: 메타데이터 딕셔너리
        """
        return {
            'name': file.name,
            'size': str(file.size),
            'type': file.type,
            'extension': os.path.splitext(file.name)[1].lower()
        }
    
    def is_supported_format(self, filename: str) -> bool:
        """
        파일 형식 지원 여부 확인
        
        Args:
            filename: 파일명
            
        Returns:
            bool: 지원 여부
        """
        extension = os.path.splitext(filename)[1].lower()
        return extension in self.supported_formats
    
    def get_file_info(self, file: UploadedFile) -> Dict[str, Union[str, int]]:
        """
        파일 정보 반환
        
        Args:
            file: Streamlit UploadedFile 객체
            
        Returns:
            Dict: 파일 정보
        """
        return {
            'name': file.name,
            'size_bytes': file.size,
            'size_kb': round(file.size / 1024, 2),
            'type': file.type,
            'extension': os.path.splitext(file.name)[1].lower(),
            'supported': self.is_supported_format(file.name)
        }
    
    def _extract_from_notebook(self, file: UploadedFile) -> str:
        """
        Jupyter Notebook에서 텍스트 내용 추출
        
        Args:
            file: Streamlit UploadedFile 객체
            
        Returns:
            str: 추출된 텍스트 내용
        """
        try:
            # 파일 내용을 문자열로 읽기
            content = file.read().decode('utf-8')
            
            # JSON 형태의 notebook 파싱
            notebook = nbformat.reads(content, as_version=4)
            
            extracted_text = []
            
            # 각 셀에서 내용 추출
            for cell in notebook.cells:
                if cell.cell_type == 'markdown':
                    # 마크다운 셀의 텍스트 추출
                    extracted_text.append(f"[Markdown Cell]\n{cell.source}\n")
                elif cell.cell_type == 'code':
                    # 코드 셀의 소스 코드 추출
                    extracted_text.append(f"[Code Cell]\n{cell.source}\n")
                    
                    # 코드 실행 결과가 있다면 추출
                    if hasattr(cell, 'outputs') and cell.outputs:
                        for output in cell.outputs:
                            if hasattr(output, 'text'):
                                extracted_text.append(f"[Output]\n{output.text}\n")
                            elif hasattr(output, 'data') and 'text/plain' in output.data:
                                extracted_text.append(f"[Output]\n{output.data['text/plain']}\n")
                elif cell.cell_type == 'raw':
                    # Raw 셀의 텍스트 추출
                    extracted_text.append(f"[Raw Cell]\n{cell.source}\n")
            
            return '\n'.join(extracted_text)
            
        except Exception as e:
            raise Exception(f"Jupyter Notebook 처리 중 오류: {str(e)}")
